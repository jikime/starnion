#!/usr/bin/env python3
"""starnion-documents — document knowledge base CLI for StarNion agent."""
import argparse, subprocess, sys, os

sys.path.insert(0, os.path.join(os.path.dirname(__file__), "..", "..", "_shared"))
from starnion_utils import _load_starnion_yaml, psql as _shared_psql

_yaml = _load_starnion_yaml()
_db   = _yaml.get("database", {}) if isinstance(_yaml.get("database"), dict) else {}
_mn   = _yaml.get("minio",    {}) if isinstance(_yaml.get("minio"),    dict) else {}

_db_url_default = (
    f"postgresql://{_db.get('user','postgres')}:{_db.get('password','')}"
    f"@{_db.get('host','localhost')}:{_db.get('port','5432')}"
    f"/{_db.get('name','starnion')}?sslmode={_db.get('ssl_mode','disable')}"
) if _db else ""

DB_URL = os.environ.get("DATABASE_URL") or _db_url_default
if not DB_URL:
    print("❌ DATABASE_URL is not set.", file=sys.stderr)
    sys.exit(1)

# MinIO — env var > YAML > default
MINIO_ENDPOINT   = os.environ.get("MINIO_ENDPOINT")   or _mn.get("endpoint",   "localhost:9000")
MINIO_ACCESS_KEY = os.environ.get("MINIO_ACCESS_KEY") or _mn.get("access_key", "")
MINIO_SECRET_KEY = os.environ.get("MINIO_SECRET_KEY") or _mn.get("secret_key", "")
MINIO_BUCKET     = os.environ.get("MINIO_BUCKET")     or _mn.get("bucket",     "starnion-files")
MINIO_USE_SSL    = (
    os.environ.get("MINIO_USE_SSL", "").lower() == "true"
    or str(_mn.get("use_ssl", "false")).lower() == "true"
)

def psql(sql): return _shared_psql(sql, DB_URL)

def esc(s):
    """Escape single quotes for SQL."""
    return (s or "").replace("'", "''")

def chunk_text(text, max_chars=1500, overlap=200):
    """Split text into overlapping chunks at natural boundaries."""
    if len(text) <= max_chars:
        return [text]
    chunks = []
    start = 0
    while start < len(text):
        end = start + max_chars
        if end < len(text):
            for sep in ["\n\n", "\n", ". ", " "]:
                pos = text.rfind(sep, start + max_chars // 2, end)
                if pos != -1:
                    end = pos + len(sep)
                    break
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start = end - overlap
    return chunks

def get_embedding(text, api_key):
    """Generate embedding via Gemini text-embedding-004. Returns list or None."""
    try:
        import google.generativeai as genai
        genai.configure(api_key=api_key)
        result = genai.embed_content(model="models/text-embedding-004", content=text)
        return result["embedding"]
    except Exception as e:
        print(f"embedding warning: {e}", file=sys.stderr)
        return None

def cmd_list(args):
    sql = (
        f"SELECT id, name, mime, size, indexed, created_at "
        f"FROM files "
        f"WHERE user_id = '{args.user_id}' AND file_type = 'document' "
        f"ORDER BY created_at DESC LIMIT 20;"
    )
    rows = psql(sql)
    if not rows:
        print("📄 No documents found.")
        return
    print("📄 Documents:")
    for row in rows.splitlines():
        parts = row.split("|")
        if len(parts) >= 6:
            doc_id, title, ftype, size, indexed, uploaded = parts
            size_kb = int(size) // 1024 if size.isdigit() else 0
            indexed_str = "indexed" if indexed == "t" else "not indexed"
            print(f"  [{doc_id}] {title} ({ftype.upper()}, {size_kb}KB) - {indexed_str} ({uploaded[:10]})")

def cmd_search(args):
    query = esc(args.query)
    limit = args.limit or 5

    # Hybrid: full-text search first, then ILIKE fallback
    sql = (
        f"SELECT fs.content, f.name, f.mime "
        f"FROM file_sections fs "
        f"JOIN files f ON f.id = fs.file_id "
        f"WHERE f.user_id = '{args.user_id}' AND f.file_type = 'document' "
        f"  AND (fs.content_tsv @@ plainto_tsquery('simple', '{query}') "
        f"       OR fs.content ILIKE '%{query}%') "
        f"ORDER BY ts_rank(fs.content_tsv, plainto_tsquery('simple', '{query}')) DESC "
        f"LIMIT {limit};"
    )
    rows = psql(sql)
    if not rows:
        print(f"🔍 No results for '{args.query}' (documents may not be indexed yet)")
        return
    print(f"🔍 Results for '{args.query}':")
    for i, row in enumerate(rows.splitlines(), 1):
        parts = row.split("|")
        if len(parts) >= 3:
            content, title, mime = parts[0], parts[1], parts[2]
            print(f"\n[{i}] Source: {title} ({mime})")
            print(f"  {content[:300]}{'...' if len(content) > 300 else ''}")

def cmd_read(args):
    sql = (
        f"SELECT fs.id, fs.content "
        f"FROM file_sections fs "
        f"JOIN files f ON f.id = fs.file_id "
        f"WHERE f.user_id = '{args.user_id}' AND f.id = {args.doc_id} "
        f"  AND f.file_type = 'document' "
        f"ORDER BY fs.id LIMIT 10;"
    )
    title_sql = (
        f"SELECT name, mime FROM files "
        f"WHERE user_id = '{args.user_id}' AND id = {args.doc_id} AND file_type = 'document';"
    )
    title_row = psql(title_sql)
    if not title_row:
        print(f"❌ Document ID {args.doc_id} not found.")
        return

    title_parts = title_row.split("|")
    doc_title = title_parts[0] if title_parts else "Unknown"
    mime = title_parts[1] if len(title_parts) > 1 else ""
    print(f"📄 Document: {doc_title} ({mime})")

    rows = psql(sql)
    if not rows:
        print("⏳ This document has not been indexed yet.")
        return

    print("\n=== Content ===")
    for row in rows.splitlines():
        parts = row.split("|")
        if len(parts) >= 2:
            content = parts[1]
            print(content)
            print()

def extract_object_key(url):
    """Extract MinIO object key from a file URL.

    Handles:
      /api/v1/files/users/{id}/files/{uuid}.ext  → users/{id}/files/{uuid}.ext
      https://host/bucket/users/{id}/...          → users/{id}/...
    Falls back to the raw URL for external files.
    """
    marker = "/api/v1/files/"
    if marker in url:
        return url.split(marker, 1)[1].split("?")[0]
    bucket = os.environ.get("MINIO_BUCKET", "starnion")
    marker2 = f"/{bucket}/"
    if marker2 in url:
        return url.split(marker2, 1)[1].split("?")[0]
    return url


def cmd_save(args):
    """Download a chat-attached file and index it into the knowledge base."""
    import urllib.request, json

    url = args.url
    filename = args.filename or url.split("/")[-1].split("?")[0] or "document"
    ext = os.path.splitext(filename)[1].lstrip(".").lower() or "bin"
    object_key = extract_object_key(url)

    try:
        with urllib.request.urlopen(url) as resp:
            data = resp.read()
    except Exception as e:
        print(f"❌ Failed to download file: {e}", file=sys.stderr)
        sys.exit(1)

    size = len(data)

    script_dir = os.path.dirname(os.path.abspath(__file__))
    extract_script = os.path.join(script_dir, "extract_text.py")
    result = subprocess.run(
        ["python3", extract_script, "--ext", ext, "--filename", filename],
        input=data, capture_output=True
    )
    if result.returncode != 0:
        print(f"❌ Text extraction failed: {result.stderr.decode().strip()}", file=sys.stderr)
        sys.exit(1)

    text = result.stdout.decode("utf-8").strip()
    if not text:
        print("❌ No text could be extracted from the file.", file=sys.stderr)
        sys.exit(1)

    doc_id = psql(
        f"INSERT INTO files (user_id, name, mime, file_type, url, object_key, size, source, indexed) "
        f"VALUES ('{esc(args.user_id)}', '{esc(filename)}', '{esc(ext)}', 'document', '{esc(url)}', "
        f"'{esc(object_key)}', {size}, 'web', false) RETURNING id"
    )
    if not doc_id:
        print("❌ Failed to create document record.", file=sys.stderr)
        sys.exit(1)

    chunks = chunk_text(text)
    gemini_key = os.environ.get("GEMINI_API_KEY", "")
    embedding_enabled = bool(gemini_key)
    indexed = 0

    for i, chunk in enumerate(chunks):
        meta = json.dumps({"chunk": i, "total": len(chunks)})
        if embedding_enabled:
            vec = get_embedding(chunk, gemini_key)
            if vec:
                vec_str = "[" + ",".join(str(v) for v in vec) + "]"
                psql(
                    f"INSERT INTO file_sections (file_id, content, embedding, content_tsv, metadata) "
                    f"VALUES ({doc_id}, '{esc(chunk)}', '{vec_str}'::vector, "
                    f"to_tsvector('simple', '{esc(chunk)}'), '{esc(meta)}'::jsonb)"
                )
            else:
                embedding_enabled = False

        if not embedding_enabled:
            psql(
                f"INSERT INTO file_sections (file_id, content, embedding, content_tsv, metadata) "
                f"VALUES ({doc_id}, '{esc(chunk)}', NULL, "
                f"to_tsvector('simple', '{esc(chunk)}'), '{esc(meta)}'::jsonb)"
            )
        indexed += 1

    psql(f"UPDATE files SET indexed = true WHERE id = {doc_id}")

    mode = "semantic+fulltext" if gemini_key else "fulltext"
    print(f"✅ '{filename}' saved to knowledge base (doc_id={doc_id}, {indexed} sections, {mode})")


# ── MinIO direct upload (AWS Signature V4, stdlib only) ──────────────────────

def _sign(key: bytes, msg: str) -> bytes:
    import hmac as _hmac, hashlib as _hashlib
    return _hmac.new(key, msg.encode("utf-8"), _hashlib.sha256).digest()


def _signing_key(secret: str, date_stamp: str) -> bytes:
    k = _sign(("AWS4" + secret).encode("utf-8"), date_stamp)
    k = _sign(k, "us-east-1")
    k = _sign(k, "s3")
    return _sign(k, "aws4_request")


def upload_to_minio(data: bytes, object_key: str, content_type: str) -> str:
    """PUT object to MinIO using AWS Signature V4. Returns the /api/files/… URL."""
    import hashlib as _hashlib, hmac as _hmac
    import urllib.request as _req, urllib.error as _uerr

    if not MINIO_ACCESS_KEY or not MINIO_SECRET_KEY:
        print(
            "❌ MinIO credentials are not set.\n"
            "   Check the minio section in ~/.starnion/starnion.yaml",
            file=sys.stderr,
        )
        sys.exit(1)

    from datetime import datetime as _dt, timezone as _tz
    scheme  = "https" if MINIO_USE_SSL else "http"
    host    = MINIO_ENDPOINT
    url     = f"{scheme}://{host}/{MINIO_BUCKET}/{object_key}"

    now        = _dt.now(_tz.utc)
    amz_date   = now.strftime("%Y%m%dT%H%M%SZ")
    date_stamp = now.strftime("%Y%m%d")

    payload_hash = _hashlib.sha256(data).hexdigest()
    canonical_headers = (
        f"content-type:{content_type}\n"
        f"host:{host}\n"
        f"x-amz-content-sha256:{payload_hash}\n"
        f"x-amz-date:{amz_date}\n"
    )
    signed_headers = "content-type;host;x-amz-content-sha256;x-amz-date"
    canonical_request = "\n".join([
        "PUT", f"/{MINIO_BUCKET}/{object_key}", "",
        canonical_headers, signed_headers, payload_hash,
    ])
    credential_scope = f"{date_stamp}/us-east-1/s3/aws4_request"
    string_to_sign = "\n".join([
        "AWS4-HMAC-SHA256", amz_date, credential_scope,
        _hashlib.sha256(canonical_request.encode()).hexdigest(),
    ])
    sig = _hmac.new(
        _signing_key(MINIO_SECRET_KEY, date_stamp),
        string_to_sign.encode(), _hashlib.sha256,
    ).hexdigest()
    authorization = (
        f"AWS4-HMAC-SHA256 Credential={MINIO_ACCESS_KEY}/{credential_scope}, "
        f"SignedHeaders={signed_headers}, Signature={sig}"
    )
    request = _req.Request(
        url, data=data,
        headers={
            "Authorization":        authorization,
            "Content-Type":         content_type,
            "x-amz-date":           amz_date,
            "x-amz-content-sha256": payload_hash,
        },
        method="PUT",
    )
    try:
        with _req.urlopen(request, timeout=60) as resp:
            resp.read()
    except _uerr.HTTPError as e:
        print(f"❌ MinIO upload error {e.code}: {e.read().decode()}", file=sys.stderr)
        sys.exit(1)
    except _uerr.URLError as e:
        print(f"❌ MinIO connection error: {e.reason}", file=sys.stderr)
        sys.exit(1)

    return f"/api/files/{object_key}"


def _index_content(doc_id: str, content: str):
    """Chunk content and insert into file_sections (with optional embeddings)."""
    import json as _json
    chunks = chunk_text(content)
    gemini_key = os.environ.get("GEMINI_API_KEY", "")
    embedding_enabled = bool(gemini_key)

    for i, chunk in enumerate(chunks):
        meta = _json.dumps({"chunk": i, "total": len(chunks)})
        if embedding_enabled:
            vec = get_embedding(chunk, gemini_key)
            if vec:
                vec_str = "[" + ",".join(str(v) for v in vec) + "]"
                psql(
                    f"INSERT INTO file_sections (file_id, content, embedding, content_tsv, metadata) "
                    f"VALUES ({doc_id}, '{esc(chunk)}', '{vec_str}'::vector, "
                    f"to_tsvector('simple', '{esc(chunk)}'), '{esc(meta)}'::jsonb)"
                )
            else:
                embedding_enabled = False
        if not embedding_enabled:
            psql(
                f"INSERT INTO file_sections (file_id, content, embedding, content_tsv, metadata) "
                f"VALUES ({doc_id}, '{esc(chunk)}', NULL, "
                f"to_tsvector('simple', '{esc(chunk)}'), '{esc(meta)}'::jsonb)"
            )

    psql(f"UPDATE files SET indexed = true WHERE id = {doc_id}")
    return len(chunks)


MIME_MAP = {
    "md":   "text/markdown",
    "txt":  "text/plain",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pdf":  "application/pdf",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "hwpx": "application/hwp+zip",
    "hwp":  "application/x-hwp",
    "html": "text/html",
}


# ── Markdown parsing (for document generation) ───────────────────────────────

import re as _re
import html as _html
from dataclasses import dataclass
from io import BytesIO as _BytesIO


# ── Korean font detection ─────────────────────────────────────────────────────

_KO_FONT_CANDIDATES = [
    # script-relative fonts/ directory
    (os.path.join(os.path.dirname(os.path.abspath(__file__)), "fonts", "NotoSansKR.ttf"), None),
    # macOS
    ("/System/Library/Fonts/AppleSDGothicNeo.ttc", 0),
    # Linux / Docker
    ("/usr/share/fonts/truetype/nanum/NanumGothic.ttf", None),
    ("/usr/share/fonts/truetype/noto/NotoSansCJK-Regular.ttc", 0),
    ("/usr/share/fonts/opentype/noto/NotoSansCJK-Regular.ttc", 0),
    ("/usr/share/fonts/noto-cjk/NotoSansCJK-Regular.ttc", 0),
]
_ko_font_name: str = ""  # set to registered font name on first use


def _register_ko_font() -> str:
    """Register a Korean-capable TTF/TTC font with reportlab, return font name."""
    global _ko_font_name
    if _ko_font_name:
        return _ko_font_name
    try:
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont
        for path, ttc_index in _KO_FONT_CANDIDATES:
            if not os.path.exists(path):
                continue
            try:
                kw = {"subfontIndex": ttc_index} if ttc_index is not None else {}
                pdfmetrics.registerFont(TTFont("KoreanFont", path, **kw))
                _ko_font_name = "KoreanFont"
                return _ko_font_name
            except Exception:
                continue
    except Exception:
        pass
    return "Helvetica"  # fallback: Korean will show as boxes


@dataclass
class _Block:
    kind: str   # h1-h6 | p | ul | ol | code_block | hr | blockquote
    text: str
    number: int = 0


@dataclass
class _Run:
    text: str
    bold: bool = False
    italic: bool = False
    code: bool = False


_INLINE_PAT = _re.compile(
    r"(\*\*\*(.+?)\*\*\*|\*\*(.+?)\*\*|__(.+?)__|"
    r"\*(.+?)\*|_(.+?)_|`([^`]+)`|\[([^\]]+)\]\([^)]+\)|!\[[^\]]*\]\([^)]+\))",
    _re.DOTALL,
)


def _parse_blocks(md: str) -> list:
    blocks, lines, i, in_code, code_buf = [], md.split("\n"), 0, False, []
    while i < len(lines):
        line = lines[i]
        if line.strip().startswith("```"):
            if in_code:
                blocks.append(_Block("code_block", "\n".join(code_buf))); code_buf = []; in_code = False
            else:
                in_code = True
            i += 1; continue
        if in_code:
            code_buf.append(line); i += 1; continue
        m = _re.match(r"^(#{1,6})\s+(.+)", line)
        if m:
            blocks.append(_Block(f"h{len(m.group(1))}", m.group(2).strip())); i += 1; continue
        if _re.match(r"^[-*_]{3,}\s*$", line):
            blocks.append(_Block("hr", "")); i += 1; continue
        m = _re.match(r"^[ \t]*[-*+]\s+(.+)", line)
        if m:
            blocks.append(_Block("ul", m.group(1))); i += 1; continue
        m = _re.match(r"^[ \t]*(\d+)[.)]\s+(.+)", line)
        if m:
            blocks.append(_Block("ol", m.group(2), number=int(m.group(1)))); i += 1; continue
        m = _re.match(r"^>\s*(.*)", line)
        if m:
            blocks.append(_Block("blockquote", m.group(1))); i += 1; continue
        if not line.strip():
            i += 1; continue
        para_lines = [line]; i += 1
        while i < len(lines):
            nxt = lines[i]
            if (not nxt.strip() or _re.match(r"^#{1,6}\s", nxt)
                    or _re.match(r"^[ \t]*[-*+]\s", nxt)
                    or _re.match(r"^[ \t]*\d+[.)]\s", nxt)
                    or _re.match(r"^[-*_]{3,}\s*$", nxt)
                    or _re.match(r"^>\s", nxt) or nxt.strip().startswith("```")):
                break
            para_lines.append(nxt); i += 1
        blocks.append(_Block("p", " ".join(para_lines)))
    return blocks


def _parse_runs(text: str) -> list:
    runs, last = [], 0
    for m in _INLINE_PAT.finditer(text):
        if m.start() > last:
            runs.append(_Run(text[last:m.start()]))
        g = m.groups()
        if g[1]:    runs.append(_Run(g[1], bold=True, italic=True))
        elif g[2]:  runs.append(_Run(g[2], bold=True))
        elif g[3]:  runs.append(_Run(g[3], bold=True))
        elif g[4]:  runs.append(_Run(g[4], italic=True))
        elif g[5]:  runs.append(_Run(g[5], italic=True))
        elif g[6]:  runs.append(_Run(g[6], code=True))
        elif g[7]:  runs.append(_Run(g[7]))
        last = m.end()
    if last < len(text):
        runs.append(_Run(text[last:]))
    return [r for r in runs if r.text]


def _plain(text: str) -> str:
    return "".join(r.text for r in _parse_runs(text))


# ── DOCX generation ───────────────────────────────────────────────────────────

def _has_cjk(text: str) -> bool:
    return any("\u1100" <= c <= "\u9fff" or "\uac00" <= c <= "\ud7af" for c in text)


_DOCX_KO_FONT  = "맑은 고딕"   # Korean body font (Word/LibreOffice built-in)
_DOCX_LAT_FONT = "Calibri"     # Latin body font
_DOCX_CODE_FONT = "Courier New"


def _set_run_font(run, text: str, code: bool = False, size_pt: int = 11) -> None:
    """Set appropriate font for a docx run, respecting CJK/Latin distinction."""
    from docx.shared import Pt
    from docx.oxml.ns import qn
    from lxml import etree

    run.font.size = Pt(size_pt)
    if code:
        if not _has_cjk(text):
            run.font.name = _DOCX_CODE_FONT
        # CJK code: leave font name unset so Word falls back to _DOCX_KO_FONT
        return

    # Set both Latin and East-Asian (rFonts) explicitly
    run.font.name = _DOCX_LAT_FONT
    rPr = run._r.get_or_add_rPr()
    rFonts = rPr.find(qn("w:rFonts"))
    if rFonts is None:
        rFonts = etree.SubElement(rPr, qn("w:rFonts"))
    rFonts.set(qn("w:eastAsia"), _DOCX_KO_FONT)
    rFonts.set(qn("w:cs"), _DOCX_KO_FONT)


def _generate_docx(title: str, content: str) -> bytes:
    """Generate DOCX from markdown using python-docx with Korean (맑은 고딕) font support."""
    from docx import Document

    doc = Document()
    doc.add_heading(title, level=1)

    for blk in _parse_blocks(content):
        if blk.kind.startswith("h") and blk.kind[1:].isdigit():
            doc.add_heading(_plain(blk.text), level=min(int(blk.kind[1:]), 9))
        elif blk.kind == "ul":
            para = doc.add_paragraph(style="List Bullet")
            for r in _parse_runs(blk.text):
                run = para.add_run(r.text)
                run.bold = r.bold; run.italic = r.italic
                _set_run_font(run, r.text, code=r.code)
        elif blk.kind == "ol":
            para = doc.add_paragraph(style="List Number")
            for r in _parse_runs(blk.text):
                run = para.add_run(r.text)
                run.bold = r.bold; run.italic = r.italic
                _set_run_font(run, r.text)
        elif blk.kind == "code_block":
            para = doc.add_paragraph()
            run = para.add_run(blk.text)
            _set_run_font(run, blk.text, code=True, size_pt=10)
        elif blk.kind == "blockquote":
            para = doc.add_paragraph()
            run = para.add_run(_plain(blk.text))
            run.italic = True
            _set_run_font(run, blk.text)
        elif blk.kind == "hr":
            doc.add_paragraph("─" * 50)
        else:
            para = doc.add_paragraph()
            for r in _parse_runs(blk.text):
                run = para.add_run(r.text)
                run.bold = r.bold; run.italic = r.italic
                _set_run_font(run, r.text, code=r.code)

    buf = _BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ── PDF generation ────────────────────────────────────────────────────────────

def _generate_pdf(title: str, content: str) -> bytes:
    """Generate PDF from markdown using reportlab with Korean font support."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.platypus import HRFlowable, Paragraph, SimpleDocTemplate, Spacer

    fn = _register_ko_font()  # "KoreanFont" if available, else "Helvetica"
    styles = getSampleStyleSheet()

    def _ps(name, **kw):
        if name not in styles:
            styles.add(ParagraphStyle(name, **kw))
        return styles[name]

    _ps("DocH1",    fontSize=20, leading=26, spaceBefore=10, spaceAfter=8,  fontName=fn)
    _ps("DocH2",    fontSize=16, leading=22, spaceBefore=8,  spaceAfter=6,  fontName=fn)
    _ps("DocH3",    fontSize=13, leading=18, spaceBefore=6,  spaceAfter=4,  fontName=fn)
    _ps("DocBody",  fontSize=11, leading=16, spaceAfter=6,                  fontName=fn)
    _ps("DocBullet",fontSize=11, leading=15, spaceAfter=4, leftIndent=14,   fontName=fn)
    _ps("DocCode",  fontSize=10, leading=14, spaceAfter=4, leftIndent=10,   fontName="Courier")
    _ps("DocQuote", fontSize=11, leading=15, spaceAfter=4, leftIndent=18,   fontName=fn)

    def _xml(text):
        parts, last = [], 0
        for m in _INLINE_PAT.finditer(text):
            if m.start() > last:
                parts.append(_html.escape(text[last:m.start()]))
            g = m.groups()
            if g[1]:   parts.append(f"<b><i>{_html.escape(g[1])}</i></b>")
            elif g[2]: parts.append(f"<b>{_html.escape(g[2])}</b>")
            elif g[3]: parts.append(f"<b>{_html.escape(g[3])}</b>")
            elif g[4]: parts.append(f"<i>{_html.escape(g[4])}</i>")
            elif g[5]: parts.append(f"<i>{_html.escape(g[5])}</i>")
            elif g[6]: parts.append(f'<font name="Courier">{_html.escape(g[6])}</font>')
            elif g[7]: parts.append(_html.escape(g[7]))
            last = m.end()
        if last < len(text):
            parts.append(_html.escape(text[last:]))
        return "".join(parts)

    buf = _BytesIO()
    doc = SimpleDocTemplate(buf, pagesize=A4,
                            topMargin=20*mm, bottomMargin=20*mm,
                            leftMargin=20*mm, rightMargin=20*mm)
    story = [Paragraph(_xml(title), styles["DocH1"]), Spacer(1, 4*mm)]

    for blk in _parse_blocks(content):
        x = _xml(blk.text)
        if blk.kind == "h1":       story.append(Paragraph(x, styles["DocH1"]))
        elif blk.kind == "h2":     story.append(Paragraph(x, styles["DocH2"]))
        elif blk.kind in ("h3","h4","h5","h6"):
                                   story.append(Paragraph(x, styles["DocH3"]))
        elif blk.kind == "ul":     story.append(Paragraph(f"• {x}", styles["DocBullet"]))
        elif blk.kind == "ol":     story.append(Paragraph(f"{blk.number}. {x}", styles["DocBullet"]))
        elif blk.kind == "code_block":
            for ln in blk.text.split("\n"):
                story.append(Paragraph(_html.escape(ln) or " ", styles["DocCode"]))
        elif blk.kind == "blockquote":
                                   story.append(Paragraph(f"│ {x}", styles["DocQuote"]))
        elif blk.kind == "hr":     story.append(HRFlowable(width="100%", thickness=0.5, spaceAfter=4))
        else:                      story.append(Paragraph(x, styles["DocBody"]))

    doc.build(story)
    return buf.getvalue()


# ── XLSX generation ───────────────────────────────────────────────────────────

def _parse_table(content: str):
    """Parse markdown table or CSV into (headers, rows)."""
    lines = [ln for ln in content.strip().split("\n") if ln.strip()]
    if not lines:
        return ["내용"], [[content]]
    if "|" in lines[0]:
        tbl = [ln for ln in lines if not _re.match(r"^\s*\|[-:| ]+\|\s*$", ln)]
        headers = [c.strip() for c in tbl[0].strip().strip("|").split("|")]
        rows = [[c.strip() for c in ln.strip().strip("|").split("|")] for ln in tbl[1:]]
        return headers, rows
    headers = [h.strip() for h in lines[0].split(",")]
    rows = [[c.strip() for c in ln.split(",")] for ln in lines[1:]]
    return headers, rows


def _generate_xlsx(title: str, content: str) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import Font
    headers, rows = _parse_table(content)
    wb = Workbook(); ws = wb.active
    ws.title = title[:31]
    ws.append(headers)
    for cell in ws[1]:
        cell.font = Font(bold=True)
    for row in rows:
        ws.append(row)
    buf = _BytesIO(); wb.save(buf); return buf.getvalue()


# ── PPTX generation ───────────────────────────────────────────────────────────

def _generate_pptx(title: str, content: str) -> bytes:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)

    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    ph_ids = {ph.placeholder_format.idx for ph in slide.placeholders}
    if 1 in ph_ids:
        slide.placeholders[1].text = ""

    bullets, pending_title = [], title

    def _flush(slide_title):
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = _plain(slide_title)
        ids = {ph.placeholder_format.idx for ph in s.placeholders}
        if 1 not in ids:
            bullets.clear(); return
        tf = s.placeholders[1].text_frame; tf.clear()
        for idx, (text, bold) in enumerate(bullets):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = text; p.font.size = Pt(18)
            if bold: p.font.bold = True
        bullets.clear()

    for blk in _parse_blocks(content):
        if blk.kind in ("h1", "h2"):
            if bullets: _flush(pending_title)
            pending_title = _plain(blk.text)
        elif blk.kind in ("h3","h4","h5","h6"):
            bullets.append((_plain(blk.text), True))
        elif blk.kind == "ul":  bullets.append(("• " + _plain(blk.text), False))
        elif blk.kind == "ol":  bullets.append((f"{blk.number}. " + _plain(blk.text), False))
        elif blk.kind == "code_block":
            for ln in blk.text.split("\n"):
                if ln.strip(): bullets.append((ln, False))
        elif blk.kind == "hr":
            if bullets: _flush(pending_title); pending_title = ""
        elif blk.kind == "blockquote": bullets.append(("❝ " + _plain(blk.text), False))
        else:  bullets.append((_plain(blk.text), False))
        if len(bullets) >= 8: _flush(pending_title)

    if bullets: _flush(pending_title)
    buf = _BytesIO(); prs.save(buf); return buf.getvalue()


# ── HWPX generation (stdlib only — no pip deps) ──────────────────────────────
# HWPX is an XML-based ZIP format (HWP Open Document Format v1.1)
# charShape IDs: 0=normal, 1=bold, 2=italic, 3=H1(20pt), 4=H2(16pt), 5=H3(13pt), 6=code
# style/paraPrShape IDs: 0=body, 1=H1, 2=H2, 3=H3, 4=list
_HWPX_HEADER_XML = """\
<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<hh:head xmlns:hh="http://www.hancom.co.kr/hwpml/2012/Head" version="1.1">
  <hh:beginNum page="1" footnote="1" endnote="1" pic="1" tbl="1" equation="1"/>
  <hh:refList>
    <hh:fontfaces>
      <hh:fontface id="0" name="함초롬바탕" type="TTF">
        <hh:font lang="HANGUL" typeface="함초롬바탕" isEmbedded="0"/>
        <hh:font lang="LATIN" typeface="함초롬바탕" isEmbedded="0"/>
        <hh:font lang="HANJA" typeface="함초롬바탕" isEmbedded="0"/>
        <hh:font lang="JAPANESE" typeface="함초롬바탕" isEmbedded="0"/>
        <hh:font lang="OTHER" typeface="함초롬바탕" isEmbedded="0"/>
        <hh:font lang="SYMBOL" typeface="함초롬바탕" isEmbedded="0"/>
        <hh:font lang="USER" typeface="함초롬바탕" isEmbedded="0"/>
      </hh:fontface>
      <hh:fontface id="1" name="Courier New" type="TTF">
        <hh:font lang="HANGUL" typeface="Courier New" isEmbedded="0"/>
        <hh:font lang="LATIN" typeface="Courier New" isEmbedded="0"/>
        <hh:font lang="HANJA" typeface="Courier New" isEmbedded="0"/>
        <hh:font lang="JAPANESE" typeface="Courier New" isEmbedded="0"/>
        <hh:font lang="OTHER" typeface="Courier New" isEmbedded="0"/>
        <hh:font lang="SYMBOL" typeface="Courier New" isEmbedded="0"/>
        <hh:font lang="USER" typeface="Courier New" isEmbedded="0"/>
      </hh:fontface>
    </hh:fontfaces>
    <hh:charShapes>
      <hh:charShape id="0" height="1000" textColor="0" shadeColor="16777215" useFontSpace="0" useKerning="0" symMark="0" borderFillIDRef="0">
        <hh:fontRef hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
        <hh:ratio hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
        <hh:spacing hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
        <hh:relSz hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
        <hh:offset hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
      </hh:charShape>
      <hh:charShape id="1" height="1000" textColor="0" shadeColor="16777215" useFontSpace="0" useKerning="0" symMark="0" borderFillIDRef="0" bold="1">
        <hh:fontRef hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
        <hh:ratio hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
        <hh:spacing hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
        <hh:relSz hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
        <hh:offset hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
      </hh:charShape>
      <hh:charShape id="2" height="1000" textColor="0" shadeColor="16777215" useFontSpace="0" useKerning="0" symMark="0" borderFillIDRef="0" italic="1">
        <hh:fontRef hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
        <hh:ratio hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
        <hh:spacing hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
        <hh:relSz hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
        <hh:offset hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
      </hh:charShape>
      <hh:charShape id="3" height="2000" textColor="0" shadeColor="16777215" useFontSpace="0" useKerning="0" symMark="0" borderFillIDRef="0" bold="1">
        <hh:fontRef hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
        <hh:ratio hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
        <hh:spacing hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
        <hh:relSz hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
        <hh:offset hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
      </hh:charShape>
      <hh:charShape id="4" height="1600" textColor="0" shadeColor="16777215" useFontSpace="0" useKerning="0" symMark="0" borderFillIDRef="0" bold="1">
        <hh:fontRef hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
        <hh:ratio hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
        <hh:spacing hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
        <hh:relSz hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
        <hh:offset hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
      </hh:charShape>
      <hh:charShape id="5" height="1300" textColor="0" shadeColor="16777215" useFontSpace="0" useKerning="0" symMark="0" borderFillIDRef="0" bold="1">
        <hh:fontRef hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
        <hh:ratio hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
        <hh:spacing hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
        <hh:relSz hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
        <hh:offset hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
      </hh:charShape>
      <hh:charShape id="6" height="1000" textColor="0" shadeColor="16777215" useFontSpace="0" useKerning="0" symMark="0" borderFillIDRef="0">
        <hh:fontRef hangul="1" latin="1" hanja="1" japanese="1" other="1" symbol="1" user="1"/>
        <hh:ratio hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
        <hh:spacing hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
        <hh:relSz hangul="100" latin="100" hanja="100" japanese="100" other="100" symbol="100" user="100"/>
        <hh:offset hangul="0" latin="0" hanja="0" japanese="0" other="0" symbol="0" user="0"/>
      </hh:charShape>
    </hh:charShapes>
    <hh:tabPrShapes>
      <hh:tabPrShape id="0"><hh:autoTabLeft/><hh:autoTabRight/></hh:tabPrShape>
    </hh:tabPrShapes>
    <hh:numberingShapes/>
    <hh:bulletShapes/>
    <hh:paraPrShapes>
      <hh:paraPrShape id="0" tabPrShapeIDRef="0" condense="0" fontLineHeight="0" snapToGrid="1" suppressLineNumbers="0" checked="0">
        <hh:align horizontal="JUSTIFY" vertical="BASELINE"/>
        <hh:heading type="NONE" idRef="0" level="0"/>
        <hh:breakSetting breakLatinWord="KEEP_WORD" breakNonLatinWord="KEEP_WORD" widowOrphan="0" keepWithNext="0" keepLines="0" pageBreakBefore="0" columnBreakBefore="0"/>
        <hh:autoSpacing eAsianEng="0" eAsianNum="0"/>
      </hh:paraPrShape>
      <hh:paraPrShape id="1" tabPrShapeIDRef="0" condense="0" fontLineHeight="0" snapToGrid="1" suppressLineNumbers="0" checked="0">
        <hh:align horizontal="LEFT" vertical="BASELINE"/>
        <hh:heading type="NONE" idRef="0" level="0"/>
        <hh:breakSetting breakLatinWord="KEEP_WORD" breakNonLatinWord="KEEP_WORD" widowOrphan="0" keepWithNext="1" keepLines="0" pageBreakBefore="0" columnBreakBefore="0"/>
        <hh:autoSpacing eAsianEng="0" eAsianNum="0"/>
      </hh:paraPrShape>
      <hh:paraPrShape id="2" tabPrShapeIDRef="0" condense="0" fontLineHeight="0" snapToGrid="1" suppressLineNumbers="0" checked="0">
        <hh:align horizontal="LEFT" vertical="BASELINE"/>
        <hh:heading type="NONE" idRef="0" level="0"/>
        <hh:breakSetting breakLatinWord="KEEP_WORD" breakNonLatinWord="KEEP_WORD" widowOrphan="0" keepWithNext="1" keepLines="0" pageBreakBefore="0" columnBreakBefore="0"/>
        <hh:autoSpacing eAsianEng="0" eAsianNum="0"/>
      </hh:paraPrShape>
      <hh:paraPrShape id="3" tabPrShapeIDRef="0" condense="0" fontLineHeight="0" snapToGrid="1" suppressLineNumbers="0" checked="0">
        <hh:align horizontal="LEFT" vertical="BASELINE"/>
        <hh:heading type="NONE" idRef="0" level="0"/>
        <hh:breakSetting breakLatinWord="KEEP_WORD" breakNonLatinWord="KEEP_WORD" widowOrphan="0" keepWithNext="1" keepLines="0" pageBreakBefore="0" columnBreakBefore="0"/>
        <hh:autoSpacing eAsianEng="0" eAsianNum="0"/>
      </hh:paraPrShape>
      <hh:paraPrShape id="4" tabPrShapeIDRef="0" condense="0" fontLineHeight="0" snapToGrid="1" suppressLineNumbers="0" checked="0">
        <hh:align horizontal="JUSTIFY" vertical="BASELINE"/>
        <hh:heading type="NONE" idRef="0" level="0"/>
        <hh:breakSetting breakLatinWord="KEEP_WORD" breakNonLatinWord="KEEP_WORD" widowOrphan="0" keepWithNext="0" keepLines="0" pageBreakBefore="0" columnBreakBefore="0"/>
        <hh:autoSpacing eAsianEng="0" eAsianNum="0"/>
      </hh:paraPrShape>
    </hh:paraPrShapes>
    <hh:styles>
      <hh:style id="0" name="바탕글" type="PARA" nextStyleIDRef="0" langID="1042" lockForm="0">
        <hh:paraPr paraPrShapeIDRef="0" tabPrShapeIDRef="0">
          <hh:paraMargin left="0" right="0" prev="0" next="0" indent="0" fixedLineHeight="0" lineHeight="160"/>
        </hh:paraPr>
        <hh:charPr charShapeIDRef="0"/>
      </hh:style>
      <hh:style id="1" name="개요 1" type="PARA" nextStyleIDRef="0" langID="1042" lockForm="0">
        <hh:paraPr paraPrShapeIDRef="1" tabPrShapeIDRef="0">
          <hh:paraMargin left="0" right="0" prev="200" next="100" indent="0" fixedLineHeight="0" lineHeight="160"/>
        </hh:paraPr>
        <hh:charPr charShapeIDRef="3"/>
      </hh:style>
      <hh:style id="2" name="개요 2" type="PARA" nextStyleIDRef="0" langID="1042" lockForm="0">
        <hh:paraPr paraPrShapeIDRef="2" tabPrShapeIDRef="0">
          <hh:paraMargin left="0" right="0" prev="150" next="80" indent="0" fixedLineHeight="0" lineHeight="160"/>
        </hh:paraPr>
        <hh:charPr charShapeIDRef="4"/>
      </hh:style>
      <hh:style id="3" name="개요 3" type="PARA" nextStyleIDRef="0" langID="1042" lockForm="0">
        <hh:paraPr paraPrShapeIDRef="3" tabPrShapeIDRef="0">
          <hh:paraMargin left="0" right="0" prev="100" next="50" indent="0" fixedLineHeight="0" lineHeight="160"/>
        </hh:paraPr>
        <hh:charPr charShapeIDRef="5"/>
      </hh:style>
      <hh:style id="4" name="목록 글" type="PARA" nextStyleIDRef="4" langID="1042" lockForm="0">
        <hh:paraPr paraPrShapeIDRef="4" tabPrShapeIDRef="0">
          <hh:paraMargin left="400" right="0" prev="0" next="0" indent="-200" fixedLineHeight="0" lineHeight="160"/>
        </hh:paraPr>
        <hh:charPr charShapeIDRef="0"/>
      </hh:style>
    </hh:styles>
    <hh:borderFills>
      <hh:borderFill id="0" threeD="0" shadow="0" centerLine="0" breakCellSeparateLine="0">
        <hh:slash close="0" countable="0"/><hh:backSlash close="0" countable="0"/>
        <hh:leftBorder type="NONE" width="0" color="0"/>
        <hh:rightBorder type="NONE" width="0" color="0"/>
        <hh:topBorder type="NONE" width="0" color="0"/>
        <hh:bottomBorder type="NONE" width="0" color="0"/>
        <hh:diagonal type="NONE" width="0" color="0"/>
        <hh:fillBrush/>
      </hh:borderFill>
    </hh:borderFills>
  </hh:refList>
  <hh:compatibleDocument targetProgram="HWP201X"><hh:layoutCompatibility/></hh:compatibleDocument>
  <hh:docOption><hh:linkInfo path="" pageInherit="1" footnoteInherit="0"/></hh:docOption>
</hh:head>"""


def _generate_hwpx(title: str, content: str) -> bytes:
    """Generate HWPX (한글 XML 포맷) from markdown — stdlib only, no pip deps."""
    import zipfile

    CS_NORMAL, CS_BOLD, CS_ITALIC, CS_H1, CS_H2, CS_H3, CS_CODE = 0, 1, 2, 3, 4, 5, 6
    ST_BODY, ST_H1, ST_H2, ST_H3, ST_LIST = 0, 1, 2, 3, 4

    pid = 0

    def para(style: int, runs: list) -> str:
        nonlocal pid
        runs_xml = "".join(
            f'<hs:run charPrIDRef="{cs}"><hs:t>{_html.escape(t)}</hs:t></hs:run>'
            for t, cs in runs if t
        )
        out = (f'<hs:p id="{pid}" paraPrIDRef="{style}" styleIDRef="{style}"'
               f' pageBreak="0" columnBreak="0">{runs_xml}</hs:p>')
        pid += 1
        return out

    def to_runs(text: str) -> list:
        result = []
        for r in _parse_runs(text):
            if r.code:            cs = CS_CODE
            elif r.bold:          cs = CS_BOLD
            elif r.italic:        cs = CS_ITALIC
            else:                 cs = CS_NORMAL
            result.append((r.text, cs))
        return result

    paragraphs = [para(ST_H1, [(title, CS_H1)])]
    for blk in _parse_blocks(content):
        k = blk.kind
        if k == "h1":
            paragraphs.append(para(ST_H1, [(_plain(blk.text), CS_H1)]))
        elif k == "h2":
            paragraphs.append(para(ST_H2, [(_plain(blk.text), CS_H2)]))
        elif k in ("h3", "h4", "h5", "h6"):
            paragraphs.append(para(ST_H3, [(_plain(blk.text), CS_H3)]))
        elif k == "ul":
            paragraphs.append(para(ST_LIST, [("• ", CS_NORMAL)] + to_runs(blk.text)))
        elif k == "ol":
            paragraphs.append(para(ST_LIST, [(f"{blk.number}. ", CS_NORMAL)] + to_runs(blk.text)))
        elif k == "code_block":
            for ln in blk.text.split("\n"):
                paragraphs.append(para(ST_BODY, [(ln or " ", CS_CODE)]))
        elif k == "hr":
            paragraphs.append(para(ST_BODY, [("─" * 40, CS_NORMAL)]))
        elif k == "blockquote":
            paragraphs.append(para(ST_BODY, [("▌ ", CS_BOLD)] + to_runs(_plain(blk.text))))
        else:
            runs = to_runs(blk.text)
            paragraphs.append(para(ST_BODY, runs or [("", CS_NORMAL)]))

    section_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        '<hs:sec xmlns:hs="http://www.hancom.co.kr/hwpml/2012/Section" version="1.1">\n'
        + "\n".join(paragraphs)
        + "\n</hs:sec>"
    )
    hpf_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        '<hp:Package xmlns:hp="http://www.hancom.co.kr/hwpml/2012/Package" version="1.1">\n'
        f'  <hp:Metadata><hp:Title>{_html.escape(title)}</hp:Title>'
        '<hp:Language>ko-KR</hp:Language>'
        '<hp:Generator>starnion-documents</hp:Generator></hp:Metadata>\n'
        '  <hp:Manifest>\n'
        '    <hp:item id="header" href="Contents/header.xml" media-type="application/hwpml-header+xml"/>\n'
        '    <hp:item id="section0" href="Contents/section0.xml" media-type="application/hwpml-section+xml"/>\n'
        '  </hp:Manifest>\n'
        '  <hp:Spine>\n'
        '    <hp:itemref idref="header"/>\n'
        '    <hp:itemref idref="section0"/>\n'
        '  </hp:Spine>\n'
        '</hp:Package>'
    )
    container_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>\n'
        '<container><rootfiles>'
        '<rootfile full-path="Contents/content.hpf" media-type="application/hwpml-package+xml"/>'
        '</rootfiles></container>'
    )
    preview = (title + "\n\n" + _re.sub(r"[#*`_~>\-]+", "", content)).strip()

    buf = _BytesIO()
    with zipfile.ZipFile(buf, "w", zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("mimetype", "application/hwp+zip")
        zf.writestr("META-INF/container.xml", container_xml)
        zf.writestr("Contents/content.hpf", hpf_xml)
        zf.writestr("Contents/header.xml", _HWPX_HEADER_XML)
        zf.writestr("Contents/section0.xml", section_xml)
        zf.writestr("Preview/PrvText.txt", preview)
    return buf.getvalue()


def _generate_hwp(title: str, content: str) -> bytes:
    """Generate HWP via LibreOffice conversion (DOCX → HWP).
    Requires: libreoffice  (apt-get install libreoffice)
    """
    import shutil, tempfile
    lo = shutil.which("libreoffice") or shutil.which("soffice")
    if not lo:
        raise RuntimeError(
            "HWP 생성에는 LibreOffice가 필요합니다.\n"
            "  Linux: apt-get install libreoffice\n"
            "  macOS: brew install --cask libreoffice\n"
            "대안: --format hwpx 를 사용하세요 (LibreOffice 불필요)"
        )
    docx_bytes = _generate_docx(title, content)
    with tempfile.TemporaryDirectory() as tmpdir:
        docx_path = os.path.join(tmpdir, "input.docx")
        with open(docx_path, "wb") as f:
            f.write(docx_bytes)
        result = subprocess.run(
            [lo, "--headless", "--convert-to", "hwp", docx_path, "--outdir", tmpdir],
            capture_output=True, timeout=60,
        )
        hwp_path = os.path.join(tmpdir, "input.hwp")
        if result.returncode == 0 and os.path.exists(hwp_path):
            with open(hwp_path, "rb") as f:
                return f.read()
        raise RuntimeError(
            f"LibreOffice 변환 실패: {result.stderr.decode(errors='replace')}"
        )


def cmd_generate(args):
    """Generate a document from AI-created content, upload to MinIO, index in DB."""
    import uuid as _uuid
    from datetime import datetime as _dt, timezone as _tz

    content = sys.stdin.read().strip()
    if not content:
        print("❌ No content received (pipe content via stdin)", file=sys.stderr)
        sys.exit(1)

    fmt = (args.format or "docx").lstrip(".")
    title = args.title
    filename = f"{title}.{fmt}"

    try:
        if fmt == "docx":
            content_bytes = _generate_docx(title, content)
        elif fmt == "pdf":
            content_bytes = _generate_pdf(title, content)
        elif fmt == "xlsx":
            content_bytes = _generate_xlsx(title, content)
        elif fmt == "pptx":
            content_bytes = _generate_pptx(title, content)
        elif fmt == "hwpx":
            content_bytes = _generate_hwpx(title, content)
        elif fmt == "hwp":
            content_bytes = _generate_hwp(title, content)
        else:
            content_bytes = content.encode("utf-8")
    except RuntimeError as e:
        print(f"❌ 문서 생성 실패 ({fmt}): {e}", file=sys.stderr)
        sys.exit(1)
    except Exception as e:
        print(f"❌ 문서 생성 실패 ({fmt}): {e}", file=sys.stderr)
        sys.exit(1)

    size = len(content_bytes)
    mime = MIME_MAP.get(fmt, "application/octet-stream")

    year = _dt.now(_tz.utc).strftime("%Y")
    object_key = f"users/{args.user_id}/docs/{year}/{_uuid.uuid4()}.{fmt}"

    print(f"Uploading to MinIO...", flush=True)
    file_url = upload_to_minio(content_bytes, object_key, mime)

    doc_id = psql(
        f"INSERT INTO files (user_id, name, mime, file_type, url, object_key, size, source, indexed) "
        f"VALUES ('{esc(args.user_id)}', '{esc(filename)}', '{esc(mime)}', 'document', '{esc(file_url)}', "
        f"'{esc(object_key)}', {size}, 'web', false) RETURNING id"
    )
    if not doc_id:
        print("❌ Failed to create document record.", file=sys.stderr)
        sys.exit(1)

    n_chunks = _index_content(doc_id, content)
    mode = "semantic+fulltext" if os.environ.get("GEMINI_API_KEY") else "fulltext"
    print(f"✅ '{filename}' generated (doc_id={doc_id}, {n_chunks} sections, {mode})")
    print(f"\n[📄 {filename}]({file_url})")


def main():
    parser = argparse.ArgumentParser(description="StarNion Document Knowledge Base")
    parser.add_argument("--user-id", required=True, help="User ID")
    sub = parser.add_subparsers(dest="cmd")

    sub.add_parser("list", help="List uploaded documents")

    p_search = sub.add_parser("search", help="Search documents")
    p_search.add_argument("--query", required=True, help="Search query")
    p_search.add_argument("--limit", type=int, default=5, help="Max results")

    p_read = sub.add_parser("read", help="Read a document's content")
    p_read.add_argument("--doc-id", required=True, type=int, help="Document ID")

    p_save = sub.add_parser("save", help="Save a chat-attached file to the knowledge base")
    p_save.add_argument("--url", required=True, help="File URL (from chat attachment)")
    p_save.add_argument("--filename", default=None, help="Original filename")

    p_gen = sub.add_parser("generate", help="Generate a document from AI content and upload to MinIO")
    p_gen.add_argument("--title", required=True, help="Document title (without extension)")
    p_gen.add_argument("--format", default="docx", choices=["docx", "pdf", "xlsx", "pptx", "hwpx", "hwp", "md", "txt"], help="File format (default: docx)")

    args = parser.parse_args()
    if not args.cmd:
        parser.print_help()
        sys.exit(1)

    dispatch = {"list": cmd_list, "search": cmd_search, "read": cmd_read, "save": cmd_save, "generate": cmd_generate}
    dispatch[args.cmd](args)

if __name__ == "__main__":
    main()
