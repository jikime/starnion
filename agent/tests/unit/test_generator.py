"""Unit tests for starnion_agent.document.generator module.

Tests cover:
- generate_pdf: PDF creation
- generate_docx: DOCX creation
- generate_xlsx: XLSX spreadsheet creation
- generate_md: Markdown file creation
- generate_txt: Plain text file creation
"""

from starnion_agent.document.generator import (
    generate_docx,
    generate_md,
    generate_pdf,
    generate_pptx,
    generate_txt,
    generate_xlsx,
)


class TestGeneratePdf:
    """Tests for generate_pdf."""

    def test_returns_bytes(self):
        """PDF output is bytes."""
        result = generate_pdf("Title", "Content")
        assert isinstance(result, bytes)

    def test_starts_with_pdf_magic(self):
        """PDF files start with %PDF magic bytes."""
        result = generate_pdf("Title", "Content")
        assert result[:5] == b"%PDF-"

    def test_nonempty_output(self):
        """Output has non-trivial size."""
        result = generate_pdf("Report", "Some content here.")
        assert len(result) > 100


class TestGenerateDocx:
    """Tests for generate_docx."""

    def test_returns_bytes(self):
        """DOCX output is bytes."""
        result = generate_docx("Title", "Content")
        assert isinstance(result, bytes)

    def test_starts_with_zip_magic(self):
        """DOCX files are ZIP archives (PK magic bytes)."""
        result = generate_docx("Title", "Content")
        assert result[:2] == b"PK"

    def test_contains_content(self):
        """Extracting text from generated DOCX recovers content (via Docling)."""
        from starnion_agent.document.parser import extract_text

        result = generate_docx("My Doc", "Hello World.\n\nSecond line.")
        text = extract_text(result, "docx", "my_doc.docx")
        assert "Hello World." in text

    def test_multiple_paragraphs(self):
        """Multiple paragraphs separated by double newlines are preserved."""
        from starnion_agent.document.parser import extract_text

        result = generate_docx("Doc", "Para 1.\n\nPara 2.\n\nPara 3.")
        text = extract_text(result, "docx", "doc.docx")
        assert "Para 1." in text
        assert "Para 2." in text
        assert "Para 3." in text


class TestGenerateXlsx:
    """Tests for generate_xlsx."""

    def test_returns_bytes(self):
        """XLSX output is bytes."""
        result = generate_xlsx(["A", "B"], [[1, 2]])
        assert isinstance(result, bytes)

    def test_starts_with_zip_magic(self):
        """XLSX files are ZIP archives."""
        result = generate_xlsx(["X"], [[1]])
        assert result[:2] == b"PK"

    def test_contains_data(self):
        """Extracting text from generated XLSX recovers data (via Docling)."""
        from starnion_agent.document.parser import extract_text

        result = generate_xlsx(
            ["Name", "Score"],
            [["Alice", 95], ["Bob", 87]],
        )
        text = extract_text(result, "xlsx", "data.xlsx")
        assert "Alice" in text
        assert "95" in text

    def test_empty_rows(self):
        """XLSX with headers only (no data rows) still works."""
        result = generate_xlsx(["Col1", "Col2"], [])
        assert isinstance(result, bytes)
        assert len(result) > 0


class TestGeneratePptx:
    """Tests for generate_pptx."""

    def test_returns_bytes(self):
        """PPTX output is bytes."""
        result = generate_pptx("Title", "Content")
        assert isinstance(result, bytes)

    def test_starts_with_zip_magic(self):
        """PPTX files are ZIP archives (PK magic bytes)."""
        result = generate_pptx("Title", "Content")
        assert result[:2] == b"PK"

    def test_nonempty_output(self):
        """Output has non-trivial size."""
        result = generate_pptx("Presentation", "Slide content here.")
        assert len(result) > 100

    def test_multiple_paragraphs(self):
        """Multiple paragraphs produce a valid PPTX."""
        content = "\n\n".join([f"Bullet {i}" for i in range(10)])
        result = generate_pptx("Many Bullets", content)
        assert isinstance(result, bytes)
        assert len(result) > 100

    def test_roundtrip_pptx(self):
        """Generated PPTX can be parsed back by python-pptx."""
        from io import BytesIO
        from pptx import Presentation

        result = generate_pptx("Test Title", "First point.\n\nSecond point.")
        prs = Presentation(BytesIO(result))
        # Title slide + 1 content slide = at least 2 slides.
        assert len(prs.slides) >= 2
        # Title slide should contain the title text.
        title_text = prs.slides[0].shapes.title.text
        assert "Test Title" in title_text

    def test_korean_content(self):
        """Korean content produces a valid PPTX."""
        result = generate_pptx("발표자료", "첫 번째 항목.\n\n두 번째 항목.")
        assert isinstance(result, bytes)
        assert result[:2] == b"PK"


class TestGenerateMd:
    """Tests for generate_md."""

    def test_returns_bytes(self):
        """Markdown output is bytes."""
        result = generate_md("Title", "Content")
        assert isinstance(result, bytes)

    def test_includes_title(self):
        """Output starts with markdown heading."""
        result = generate_md("My Report", "Some content.")
        text = result.decode("utf-8")
        assert text.startswith("# My Report")

    def test_includes_content(self):
        """Content follows the title."""
        result = generate_md("Title", "Body text here.")
        text = result.decode("utf-8")
        assert "Body text here." in text

    def test_utf8_encoding(self):
        """Korean content is properly UTF-8 encoded."""
        result = generate_md("제목", "내용입니다.")
        text = result.decode("utf-8")
        assert "제목" in text
        assert "내용입니다." in text


class TestGenerateTxt:
    """Tests for generate_txt."""

    def test_returns_bytes(self):
        """Text output is bytes."""
        result = generate_txt("Hello")
        assert isinstance(result, bytes)

    def test_content_preserved(self):
        """Content is exactly preserved."""
        result = generate_txt("Line 1\nLine 2")
        assert result == b"Line 1\nLine 2"

    def test_empty_content(self):
        """Empty content returns empty bytes."""
        result = generate_txt("")
        assert result == b""
