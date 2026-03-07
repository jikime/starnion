"""Document parsing: extract text from various document formats."""

import logging
import re
import struct
from collections.abc import Callable
from io import BytesIO

import httpx
from pypdf import PdfReader

logger = logging.getLogger(__name__)

_DOC_FALLBACK = "(DOC 파일에서 텍스트를 추출할 수 없었어요. .docx 포맷으로 변환해서 보내주세요.)"
_XLS_FALLBACK = "(XLS 파일에서 텍스트를 추출할 수 없었어요. .xlsx 포맷으로 변환해서 보내주세요.)"
_PPT_FALLBACK = "(PPT 파일에서 텍스트를 추출할 수 없었어요. .pptx 포맷으로 변환해서 보내주세요.)"


async def fetch_file(url: str) -> bytes:
    """Download a file from *url* and return its bytes."""
    async with httpx.AsyncClient(timeout=60) as client:
        resp = await client.get(url)
        resp.raise_for_status()
        return resp.content


def extract_text_from_pdf(data: bytes) -> str:
    """Extract text from all pages of a PDF byte-string."""
    reader = PdfReader(BytesIO(data))
    pages: list[str] = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text.strip())
    return "\n\n".join(pages)


def extract_text_from_docx(data: bytes) -> str:
    """Extract text from a DOCX file."""
    from docx import Document

    doc = Document(BytesIO(data))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def extract_text_from_xlsx(data: bytes) -> str:
    """Extract text from an XLSX spreadsheet (all sheets)."""
    from openpyxl import load_workbook

    wb = load_workbook(BytesIO(data), read_only=True, data_only=True)
    lines: list[str] = []
    for ws in wb.worksheets:
        lines.append(f"## {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                lines.append("\t".join(cells))
    wb.close()
    return "\n".join(lines)


def extract_text_from_pptx(data: bytes) -> str:
    """Extract text from a PPTX presentation."""
    from pptx import Presentation

    prs = Presentation(BytesIO(data))
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            slides.append(f"## Slide {i}\n" + "\n".join(texts))
    return "\n\n".join(slides)


def extract_text_from_doc(data: bytes) -> str:
    """Extract text from a legacy DOC (Word Binary) file.

    Strategy:
    1. Try python-docx (handles .docx files saved with .doc extension)
    2. Try olefile + Word Binary Format FIB/piece-table parsing
    3. Fallback with conversion guidance
    """
    # Strategy 1: python-docx for mislabeled .docx files
    try:
        from docx import Document

        doc = Document(BytesIO(data))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        if paragraphs:
            return "\n\n".join(paragraphs)
    except Exception:
        pass

    # Strategy 2: olefile-based Word Binary Format parsing
    try:
        text = _parse_doc_binary(data)
        if text:
            return text
    except Exception:
        logger.debug("DOC binary parsing failed", exc_info=True)

    return _DOC_FALLBACK


def _parse_doc_binary(data: bytes) -> str | None:
    """Parse Word Binary Format (97-2003) to extract document text.

    Reads the FIB (File Information Block) to locate the CLX (piece table)
    in the Table stream, then extracts text pieces (ANSI or Unicode).
    """
    import olefile

    buf = BytesIO(data)
    if not olefile.isOleFile(buf):
        return None

    ole = olefile.OleFileIO(buf)
    try:
        if not ole.exists("WordDocument"):
            return None

        ws = ole.openstream("WordDocument").read()

        # Validate Word Binary magic number (0xA5EC)
        if len(ws) < 0x60 or struct.unpack_from("<H", ws, 0)[0] != 0xA5EC:
            return None

        # Determine Table stream (bit 9 of flags → 0Table or 1Table)
        flags = struct.unpack_from("<H", ws, 0x0A)[0]
        tbl_name = "1Table" if (flags & 0x0200) else "0Table"
        if not ole.exists(tbl_name):
            return None

        # Parse FIB offsets dynamically:
        #   FibBase(32) → csw → FibRgW → cslw → FibRgLw → cbRgFcLcb → FibRgFcLcb
        csw = struct.unpack_from("<H", ws, 0x0020)[0]
        off = 0x0022 + csw * 2  # end of FibRgW
        cslw = struct.unpack_from("<H", ws, off)[0]
        off += 2 + cslw * 4  # end of FibRgLw
        cb_fc = struct.unpack_from("<H", ws, off)[0]  # cbRgFcLcb
        fc_start = off + 2  # start of FibRgFcLcb

        # fcClx/lcbClx = pair #33 in FibRgFcLcb (MS-DOC spec §2.5.10)
        _CLX_PAIR = 33
        if cb_fc <= _CLX_PAIR:
            return None

        fc_clx = struct.unpack_from("<I", ws, fc_start + _CLX_PAIR * 8)[0]
        lcb_clx = struct.unpack_from("<I", ws, fc_start + _CLX_PAIR * 8 + 4)[0]
        if lcb_clx == 0:
            return None

        # Read CLX from Table stream
        ts = ole.openstream(tbl_name).read()
        if fc_clx + lcb_clx > len(ts):
            return None
        clx = ts[fc_clx : fc_clx + lcb_clx]

        # Skip PRC entries (type 0x01) to reach Pcdt (type 0x02)
        pos = 0
        while pos < len(clx) and clx[pos] == 0x01:
            cb_prc = struct.unpack_from("<H", clx, pos + 1)[0]
            pos += 3 + cb_prc
        if pos >= len(clx) or clx[pos] != 0x02:
            return None

        lcb_pcdt = struct.unpack_from("<I", clx, pos + 1)[0]
        pcd = clx[pos + 5 : pos + 5 + lcb_pcdt]

        # Parse PlcPcd: (n+1) CPs (4 bytes each) + n PCDs (8 bytes each)
        n = (lcb_pcdt - 4) // 12
        if n <= 0:
            return None

        parts: list[str] = []
        for i in range(n):
            cp0 = struct.unpack_from("<I", pcd, i * 4)[0]
            cp1 = struct.unpack_from("<I", pcd, (i + 1) * 4)[0]
            pcd_off = (n + 1) * 4 + i * 8
            if pcd_off + 8 > len(pcd):
                break

            fc_val = struct.unpack_from("<I", pcd, pcd_off + 2)[0]
            is_ansi = bool(fc_val & 0x40000000)
            fc_real = fc_val & 0x3FFFFFFF
            count = cp1 - cp0
            if count <= 0:
                continue

            try:
                if is_ansi:
                    start = fc_real // 2
                    chunk = ws[start : start + count]
                    # CP949 for Korean, falls back gracefully for Western text
                    parts.append(chunk.decode("cp949", errors="replace"))
                else:
                    chunk = ws[fc_real : fc_real + count * 2]
                    parts.append(chunk.decode("utf-16-le", errors="replace"))
            except (IndexError, struct.error):
                continue

        if not parts:
            return None

        text = "".join(parts)
        text = text.replace("\r\n", "\n").replace("\r", "\n")
        text = re.sub(r"[\x00-\x08\x0b\x0c\x0e-\x1f]", "", text)
        return text.strip() or None
    finally:
        ole.close()


def extract_text_from_xls(data: bytes) -> str:
    """Extract text from a legacy XLS file.

    Tries openpyxl first (for mislabeled .xlsx), then returns fallback.
    """
    try:
        return extract_text_from_xlsx(data)
    except Exception:
        pass
    return _XLS_FALLBACK


def extract_text_from_ppt(data: bytes) -> str:
    """Extract text from a legacy PPT file.

    Tries python-pptx first (for mislabeled .pptx), then returns fallback.
    """
    try:
        return extract_text_from_pptx(data)
    except Exception:
        pass
    return _PPT_FALLBACK


def extract_text_from_hwp(data: bytes) -> str:
    """Extract text from a HWP file (best-effort via olefile)."""
    try:
        import olefile

        ole = olefile.OleFileIO(BytesIO(data))
        if ole.exists("PrvText"):
            raw = ole.openstream("PrvText").read()
            text = raw.decode("utf-16-le", errors="replace")
            ole.close()
            return text.strip()
        ole.close()
        return "(HWP 파일에서 텍스트를 추출할 수 없었어요. PrvText 스트림이 없습니다.)"
    except Exception:
        logger.debug("HWP extraction failed", exc_info=True)
        return "(HWP 파일 처리 중 오류가 발생했어요. 다른 포맷으로 변환해서 보내주세요.)"


def extract_text_from_md(data: bytes) -> str:
    """Extract text from a Markdown file."""
    return data.decode("utf-8", errors="replace").strip()


def extract_text_from_txt(data: bytes) -> str:
    """Extract text from a plain text file."""
    return data.decode("utf-8", errors="replace").strip()


# Extension → extractor mapping.
_EXTRACTORS: dict[str, Callable[[bytes], str]] = {
    "pdf": extract_text_from_pdf,
    "docx": extract_text_from_docx,
    "doc": extract_text_from_doc,
    "xlsx": extract_text_from_xlsx,
    "xls": extract_text_from_xls,
    "pptx": extract_text_from_pptx,
    "ppt": extract_text_from_ppt,
    "hwp": extract_text_from_hwp,
    "md": extract_text_from_md,
    "markdown": extract_text_from_md,
    "txt": extract_text_from_txt,
    "text": extract_text_from_txt,
    "csv": extract_text_from_txt,
}


def extract_text(data: bytes, ext: str) -> str:
    """Route to the appropriate extractor based on file extension.

    Args:
        data: Raw file bytes.
        ext: File extension without dot (e.g. "pdf", "docx").

    Returns:
        Extracted text or an error message.
    """
    extractor = _EXTRACTORS.get(ext.lower())
    if extractor is None:
        return f"(지원하지 않는 파일 형식이에요: .{ext})"
    return extractor(data)
