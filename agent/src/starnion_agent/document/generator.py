"""Document generation: create documents in various formats.

All generators accept raw LLM markdown and convert it to proper formatting.
"""

from __future__ import annotations

import html
import re
from dataclasses import dataclass
from io import BytesIO
from pathlib import Path

# ── Font setup ────────────────────────────────────────────────────────────────

_FONT_DIR = Path(__file__).parent / "fonts"
_NOTO_SANS_KR = _FONT_DIR / "NotoSansKR.ttf"
_FONT_REGISTERED = False


def _ensure_font() -> str:
    global _FONT_REGISTERED  # noqa: PLW0603
    if not _FONT_REGISTERED and _NOTO_SANS_KR.exists():
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont

        pdfmetrics.registerFont(TTFont("NotoSansKR", str(_NOTO_SANS_KR)))
        _FONT_REGISTERED = True
    return "NotoSansKR" if _FONT_REGISTERED else "Helvetica"


# ── Markdown parsing ──────────────────────────────────────────────────────────

@dataclass
class Block:
    """A structural block parsed from markdown."""
    kind: str   # h1-h6 | p | ul | ol | code_block | hr | blockquote
    text: str
    number: int = 0  # for ol items


@dataclass
class Run:
    """An inline run with formatting."""
    text: str
    bold: bool = False
    italic: bool = False
    code: bool = False


def _parse_blocks(md: str) -> list[Block]:
    """Parse markdown into a list of structural blocks."""
    blocks: list[Block] = []
    lines = md.split("\n")
    i = 0
    in_code = False
    code_buf: list[str] = []

    while i < len(lines):
        line = lines[i]

        # Fenced code block
        if line.strip().startswith("```"):
            if in_code:
                blocks.append(Block("code_block", "\n".join(code_buf)))
                code_buf = []
                in_code = False
            else:
                in_code = True
            i += 1
            continue

        if in_code:
            code_buf.append(line)
            i += 1
            continue

        # Heading
        m = re.match(r"^(#{1,6})\s+(.+)", line)
        if m:
            level = len(m.group(1))
            blocks.append(Block(f"h{level}", m.group(2).strip()))
            i += 1
            continue

        # Horizontal rule
        if re.match(r"^[-*_]{3,}\s*$", line):
            blocks.append(Block("hr", ""))
            i += 1
            continue

        # Unordered list
        m = re.match(r"^[ \t]*[-*+]\s+(.+)", line)
        if m:
            blocks.append(Block("ul", m.group(1)))
            i += 1
            continue

        # Ordered list
        m = re.match(r"^[ \t]*(\d+)[.)]\s+(.+)", line)
        if m:
            blocks.append(Block("ol", m.group(2), number=int(m.group(1))))
            i += 1
            continue

        # Blockquote
        m = re.match(r"^>\s*(.*)", line)
        if m:
            blocks.append(Block("blockquote", m.group(1)))
            i += 1
            continue

        # Blank line — skip
        if not line.strip():
            i += 1
            continue

        # Paragraph: accumulate until blank / heading / list
        para_lines = [line]
        i += 1
        while i < len(lines):
            nxt = lines[i]
            if (not nxt.strip()
                    or re.match(r"^#{1,6}\s", nxt)
                    or re.match(r"^[ \t]*[-*+]\s", nxt)
                    or re.match(r"^[ \t]*\d+[.)]\s", nxt)
                    or re.match(r"^[-*_]{3,}\s*$", nxt)
                    or re.match(r"^>\s", nxt)
                    or nxt.strip().startswith("```")):
                break
            para_lines.append(nxt)
            i += 1
        blocks.append(Block("p", " ".join(para_lines)))

    return blocks


_INLINE_PATTERN = re.compile(
    r"(\*\*\*(.+?)\*\*\*"      # bold+italic
    r"|\*\*(.+?)\*\*"           # bold **
    r"|__(.+?)__"               # bold __
    r"|\*(.+?)\*"               # italic *
    r"|_(.+?)_"                 # italic _
    r"|`([^`]+)`"               # inline code
    r"|\[([^\]]+)\]\([^)]+\)"   # link [text](url)
    r"|!\[[^\]]*\]\([^)]+\))",  # image — drop
    re.DOTALL,
)


def _parse_runs(text: str) -> list[Run]:
    """Parse inline markdown into a list of formatted runs."""
    runs: list[Run] = []
    last = 0
    for m in _INLINE_PATTERN.finditer(text):
        if m.start() > last:
            runs.append(Run(text[last:m.start()]))
        if m.group(2):          # ***bold+italic***
            runs.append(Run(m.group(2), bold=True, italic=True))
        elif m.group(3):        # **bold**
            runs.append(Run(m.group(3), bold=True))
        elif m.group(4):        # __bold__
            runs.append(Run(m.group(4), bold=True))
        elif m.group(5):        # *italic*
            runs.append(Run(m.group(5), italic=True))
        elif m.group(6):        # _italic_
            runs.append(Run(m.group(6), italic=True))
        elif m.group(7):        # `code`
            runs.append(Run(m.group(7), code=True))
        elif m.group(8):        # [link](url) → link text only
            runs.append(Run(m.group(8)))
        # image: skip
        last = m.end()
    if last < len(text):
        runs.append(Run(text[last:]))
    return [r for r in runs if r.text]


def _plain(text: str) -> str:
    """Strip all inline markdown to plain text."""
    return "".join(r.text for r in _parse_runs(text))


def _to_reportlab_xml(text: str) -> str:
    """Convert inline markdown to ReportLab XML markup string."""
    # Escape HTML special chars in raw text first, then mark up.
    # We build output segment by segment to avoid double-escaping the tags.
    parts: list[str] = []
    last = 0
    for m in _INLINE_PATTERN.finditer(text):
        # Plain text before match — escape
        if m.start() > last:
            parts.append(html.escape(text[last:m.start()]))
        if m.group(2):
            parts.append(f"<b><i>{html.escape(m.group(2))}</i></b>")
        elif m.group(3):
            parts.append(f"<b>{html.escape(m.group(3))}</b>")
        elif m.group(4):
            parts.append(f"<b>{html.escape(m.group(4))}</b>")
        elif m.group(5):
            parts.append(f"<i>{html.escape(m.group(5))}</i>")
        elif m.group(6):
            parts.append(f"<i>{html.escape(m.group(6))}</i>")
        elif m.group(7):
            parts.append(f'<font name="Courier">{html.escape(m.group(7))}</font>')
        elif m.group(8):
            parts.append(html.escape(m.group(8)))
        # image: drop
        last = m.end()
    if last < len(text):
        parts.append(html.escape(text[last:]))
    return "".join(parts)


# ── PDF ───────────────────────────────────────────────────────────────────────

def generate_pdf(title: str, content: str) -> bytes:
    """Generate a PDF document with full Korean/CJK support, rendering markdown."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.platypus import HRFlowable, Paragraph, SimpleDocTemplate, Spacer

    fn = _ensure_font()
    styles = getSampleStyleSheet()

    def _ps(name: str, **kw):  # type: ignore[return]
        if name not in styles:
            styles.add(ParagraphStyle(name, **kw))
        return styles[name]

    _ps("KorH1",    fontSize=20, leading=26, spaceBefore=10, spaceAfter=8, fontName=fn)
    _ps("KorH2",    fontSize=16, leading=22, spaceBefore=8,  spaceAfter=6, fontName=fn)
    _ps("KorH3",    fontSize=13, leading=18, spaceBefore=6,  spaceAfter=4, fontName=fn)
    _ps("KorH4",    fontSize=12, leading=16, spaceBefore=4,  spaceAfter=3, fontName=fn)
    _ps("KorBody",  fontSize=11, leading=16, spaceAfter=6,   fontName=fn)
    _ps("KorBullet",fontSize=11, leading=15, spaceAfter=4,   leftIndent=14, fontName=fn)
    _ps("KorCode",  fontSize=10, leading=14, spaceAfter=4,   leftIndent=10,
        fontName="Courier", backColor="#f4f4f4")
    _ps("KorQuote", fontSize=11, leading=15, spaceAfter=4,   leftIndent=18,
        textColor="#555555", fontName=fn)

    buf = BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        topMargin=20 * mm, bottomMargin=20 * mm,
        leftMargin=20 * mm, rightMargin=20 * mm,
    )

    story: list = [Paragraph(_to_reportlab_xml(title), styles["KorH1"]), Spacer(1, 4 * mm)]

    for blk in _parse_blocks(content):
        xml = _to_reportlab_xml(blk.text)
        if blk.kind == "h1":
            story.append(Paragraph(xml, styles["KorH1"]))
        elif blk.kind == "h2":
            story.append(Paragraph(xml, styles["KorH2"]))
        elif blk.kind in ("h3", "h4", "h5", "h6"):
            story.append(Paragraph(xml, styles["KorH3"] if blk.kind == "h3" else styles["KorH4"]))
        elif blk.kind == "ul":
            story.append(Paragraph(f"• {xml}", styles["KorBullet"]))
        elif blk.kind == "ol":
            story.append(Paragraph(f"{blk.number}. {xml}", styles["KorBullet"]))
        elif blk.kind == "code_block":
            for ln in blk.text.split("\n"):
                story.append(Paragraph(html.escape(ln) or " ", styles["KorCode"]))
        elif blk.kind == "blockquote":
            story.append(Paragraph(f"│ {xml}", styles["KorQuote"]))
        elif blk.kind == "hr":
            story.append(HRFlowable(width="100%", thickness=0.5, color="#cccccc", spaceAfter=4))
        else:  # p
            story.append(Paragraph(xml, styles["KorBody"]))

    doc.build(story)
    return buf.getvalue()


# ── DOCX ─────────────────────────────────────────────────────────────────────

def generate_docx(title: str, content: str) -> bytes:
    """Generate a DOCX document from markdown content."""
    from docx import Document
    from docx.shared import Pt, RGBColor

    doc = Document()
    doc.add_heading(title, level=1)

    for blk in _parse_blocks(content):
        if blk.kind.startswith("h") and blk.kind[1:].isdigit():
            level = int(blk.kind[1:])
            doc.add_heading(_plain(blk.text), level=min(level, 9))

        elif blk.kind == "ul":
            para = doc.add_paragraph(style="List Bullet")
            _add_runs_to_para(para, _parse_runs(blk.text))

        elif blk.kind == "ol":
            para = doc.add_paragraph(style="List Number")
            _add_runs_to_para(para, _parse_runs(blk.text))

        elif blk.kind == "code_block":
            para = doc.add_paragraph()
            run = para.add_run(blk.text)
            run.font.name = "Courier New"
            run.font.size = Pt(10)

        elif blk.kind == "blockquote":
            para = doc.add_paragraph()
            run = para.add_run(_plain(blk.text))
            run.font.italic = True
            run.font.color.rgb = RGBColor(0x55, 0x55, 0x55)

        elif blk.kind == "hr":
            doc.add_paragraph("─" * 50)

        else:  # p
            para = doc.add_paragraph()
            _add_runs_to_para(para, _parse_runs(blk.text))

    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def _add_runs_to_para(para, runs: list[Run]) -> None:
    """Add formatted runs to a python-docx paragraph."""
    from docx.shared import Pt

    for r in runs:
        run = para.add_run(r.text)
        run.bold = r.bold
        run.italic = r.italic
        if r.code:
            run.font.name = "Courier New"
            run.font.size = Pt(10)


# ── XLSX ─────────────────────────────────────────────────────────────────────

def generate_xlsx(headers: list[str], rows: list[list]) -> bytes:
    """Generate an XLSX spreadsheet from tabular data."""
    from openpyxl import Workbook
    from openpyxl.styles import Font

    wb = Workbook()
    ws = wb.active
    assert ws is not None
    ws.append([_plain(h) for h in headers])
    # Bold header row
    for cell in ws[1]:
        cell.font = Font(bold=True)
    for row in rows:
        ws.append([_plain(str(c)) for c in row])
    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


def parse_xlsx_content(content: str) -> tuple[list[str], list[list[str]]]:
    """Parse LLM content (markdown table or CSV) into headers + rows."""
    lines = [ln for ln in content.strip().split("\n") if ln.strip()]
    if not lines:
        return ["내용"], [[content]]

    # Detect markdown table: lines contain "|"
    if "|" in lines[0]:
        table_lines = [ln for ln in lines if not re.match(r"^\s*\|[-:| ]+\|\s*$", ln)]
        headers = [c.strip() for c in table_lines[0].strip().strip("|").split("|")]
        rows = [
            [c.strip() for c in ln.strip().strip("|").split("|")]
            for ln in table_lines[1:]
        ]
        return headers, rows

    # CSV fallback
    headers = [h.strip() for h in lines[0].split(",")]
    rows = [[c.strip() for c in ln.split(",")] for ln in lines[1:]]
    return headers, rows


# ── PPTX ─────────────────────────────────────────────────────────────────────

def generate_pptx(title: str, content: str) -> bytes:
    """Generate a PPTX presentation from markdown content.

    Strategy:
    - h1 / h2 headings → new slide (used as slide title)
    - ul / ol / p      → bullet items on the current slide
    - code_block       → bullet with monospace hint
    - hr               → new blank slide
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    # Title slide
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    if 1 in {ph.placeholder_format.idx for ph in slide.placeholders}:
        slide.placeholders[1].text = ""

    current_slide = None
    current_tf = None
    bullets: list[tuple[str, bool]] = []  # (text, is_bold)

    def _flush(slide_title: str) -> None:
        nonlocal current_slide, current_tf
        current_slide = prs.slides.add_slide(content_layout)
        current_slide.shapes.title.text = _plain(slide_title)
        ph_ids = {ph.placeholder_format.idx for ph in current_slide.placeholders}
        if 1 not in ph_ids:
            bullets.clear()
            current_tf = None
            return
        tf = current_slide.placeholders[1].text_frame
        current_tf = tf
        tf.clear()
        first = True
        for text, bold in bullets:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.text = text
            p.font.size = Pt(18)
            if bold:
                p.font.bold = True
        bullets.clear()

    pending_title = title
    for blk in _parse_blocks(content):
        if blk.kind in ("h1", "h2"):
            if bullets:
                _flush(pending_title)
            pending_title = _plain(blk.text)
        elif blk.kind in ("h3", "h4", "h5", "h6"):
            bullets.append((_plain(blk.text), True))
        elif blk.kind == "ul":
            bullets.append(("• " + _plain(blk.text), False))
        elif blk.kind == "ol":
            bullets.append((f"{blk.number}. {_plain(blk.text)}", False))
        elif blk.kind == "code_block":
            for ln in blk.text.split("\n"):
                if ln.strip():
                    bullets.append((ln, False))
        elif blk.kind == "hr":
            if bullets:
                _flush(pending_title)
            pending_title = ""
        elif blk.kind == "blockquote":
            bullets.append(("❝ " + _plain(blk.text), False))
        else:  # p
            bullets.append((_plain(blk.text), False))

        # Split slide when too many bullets
        if len(bullets) >= 8:
            _flush(pending_title)

    if bullets:
        _flush(pending_title)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Markdown ─────────────────────────────────────────────────────────────────

def generate_md(title: str, content: str) -> bytes:
    """Generate a Markdown document (content is already markdown — pass through)."""
    text = f"# {title}\n\n{content}"
    return text.encode("utf-8")


# ── TXT ──────────────────────────────────────────────────────────────────────

def generate_txt(content: str) -> bytes:
    """Generate plain text by stripping inline markdown syntax, preserving line structure."""
    lines = [_plain(line) for line in content.split("\n")]
    return "\n".join(lines).encode("utf-8")
