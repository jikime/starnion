"""Document generation: create documents in various formats."""

from io import BytesIO
from pathlib import Path

# Bundled Noto Sans KR font for Korean/CJK PDF support.
_FONT_DIR = Path(__file__).parent / "fonts"
_NOTO_SANS_KR = _FONT_DIR / "NotoSansKR.ttf"

# ReportLab font registration (module-level, runs once).
_FONT_REGISTERED = False


def _ensure_font() -> str:
    """Register Noto Sans KR with ReportLab and return the font name."""
    global _FONT_REGISTERED  # noqa: PLW0603
    if not _FONT_REGISTERED and _NOTO_SANS_KR.exists():
        from reportlab.pdfbase import pdfmetrics
        from reportlab.pdfbase.ttfonts import TTFont

        pdfmetrics.registerFont(TTFont("NotoSansKR", str(_NOTO_SANS_KR)))
        _FONT_REGISTERED = True
    return "NotoSansKR" if _FONT_REGISTERED else "Helvetica"


def generate_pdf(title: str, content: str) -> bytes:
    """Generate a PDF document with full Korean/CJK support using ReportLab."""
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
    from reportlab.lib.units import mm
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    font_name = _ensure_font()
    styles = getSampleStyleSheet()
    styles.add(ParagraphStyle(
        "KorTitle", fontName=font_name, fontSize=18, leading=24, spaceAfter=12,
    ))
    styles.add(ParagraphStyle(
        "KorBody", fontName=font_name, fontSize=11, leading=16, spaceAfter=8,
    ))

    buf = BytesIO()
    doc = SimpleDocTemplate(
        buf, pagesize=A4,
        topMargin=20 * mm, bottomMargin=20 * mm,
        leftMargin=20 * mm, rightMargin=20 * mm,
    )

    story: list = [Paragraph(title, styles["KorTitle"]), Spacer(1, 6 * mm)]
    for para in content.split("\n\n"):
        text = para.strip()
        if text:
            # Replace single newlines with <br/> for ReportLab.
            story.append(Paragraph(text.replace("\n", "<br/>"), styles["KorBody"]))

    doc.build(story)
    return buf.getvalue()


def generate_docx(title: str, content: str) -> bytes:
    """Generate a DOCX document from text content."""
    from docx import Document

    doc = Document()
    doc.add_heading(title, level=1)
    for paragraph in content.split("\n\n"):
        if paragraph.strip():
            doc.add_paragraph(paragraph.strip())
    buf = BytesIO()
    doc.save(buf)
    return buf.getvalue()


def generate_xlsx(headers: list[str], rows: list[list]) -> bytes:
    """Generate an XLSX spreadsheet from tabular data."""
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.append(headers)
    for row in rows:
        ws.append(row)
    buf = BytesIO()
    wb.save(buf)
    return buf.getvalue()


def generate_pptx(title: str, content: str) -> bytes:
    """Generate a PPTX presentation from text content.

    The first slide is a title slide.  Each ``\\n\\n``-delimited paragraph
    becomes a bullet on a content slide (new slide every 6 bullets).
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide.
    title_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_layout)
    slide.shapes.title.text = title
    if 1 in slide.placeholders:
        slide.placeholders[1].text = ""

    # Content slides from paragraphs.
    paragraphs = [p.strip() for p in content.split("\n\n") if p.strip()]
    bullets_per_slide = 6
    for i in range(0, len(paragraphs), bullets_per_slide):
        chunk = paragraphs[i:i + bullets_per_slide]
        content_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = title
        if 1 not in slide.placeholders:
            continue
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        for j, para_text in enumerate(chunk):
            if j == 0:
                tf.paragraphs[0].text = para_text
                tf.paragraphs[0].font.size = Pt(18)
            else:
                p = tf.add_paragraph()
                p.text = para_text
                p.font.size = Pt(18)

    buf = BytesIO()
    prs.save(buf)
    return buf.getvalue()


def generate_md(title: str, content: str) -> bytes:
    """Generate a Markdown document."""
    text = f"# {title}\n\n{content}"
    return text.encode("utf-8")


def generate_txt(content: str) -> bytes:
    """Generate a plain text document."""
    return content.encode("utf-8")
