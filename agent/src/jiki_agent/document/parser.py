"""Document parsing: extract text from various document formats."""

import logging
from collections.abc import Callable
from io import BytesIO

import httpx
from pypdf import PdfReader

logger = logging.getLogger(__name__)


async def fetch_file(url: str) -> bytes:
    """Download a file from *url* and return its bytes."""
    async with httpx.AsyncClient(timeout=60) as client:
        resp = await client.get(url)
        resp.raise_for_status()
        return resp.content


def extract_text_from_pdf(data: bytes) -> str:
    """Extract text from all pages of a PDF byte-string."""
    reader = PdfReader(BytesIO(data))
    pages: list[str] = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text.strip())
    return "\n\n".join(pages)


def extract_text_from_docx(data: bytes) -> str:
    """Extract text from a DOCX file."""
    from docx import Document

    doc = Document(BytesIO(data))
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n\n".join(paragraphs)


def extract_text_from_xlsx(data: bytes) -> str:
    """Extract text from an XLSX spreadsheet (all sheets)."""
    from openpyxl import load_workbook

    wb = load_workbook(BytesIO(data), read_only=True, data_only=True)
    lines: list[str] = []
    for ws in wb.worksheets:
        lines.append(f"## {ws.title}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                lines.append("\t".join(cells))
    wb.close()
    return "\n".join(lines)


def extract_text_from_pptx(data: bytes) -> str:
    """Extract text from a PPTX presentation."""
    from pptx import Presentation

    prs = Presentation(BytesIO(data))
    slides: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        texts: list[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
        if texts:
            slides.append(f"## Slide {i}\n" + "\n".join(texts))
    return "\n\n".join(slides)


def extract_text_from_hwp(data: bytes) -> str:
    """Extract text from a HWP file (best-effort via olefile)."""
    try:
        import olefile

        ole = olefile.OleFileIO(BytesIO(data))
        if ole.exists("PrvText"):
            raw = ole.openstream("PrvText").read()
            text = raw.decode("utf-16-le", errors="replace")
            ole.close()
            return text.strip()
        ole.close()
        return "(HWP 파일에서 텍스트를 추출할 수 없었어요. PrvText 스트림이 없습니다.)"
    except Exception:
        logger.debug("HWP extraction failed", exc_info=True)
        return "(HWP 파일 처리 중 오류가 발생했어요. 다른 포맷으로 변환해서 보내주세요.)"


def extract_text_from_md(data: bytes) -> str:
    """Extract text from a Markdown file."""
    return data.decode("utf-8", errors="replace").strip()


def extract_text_from_txt(data: bytes) -> str:
    """Extract text from a plain text file."""
    return data.decode("utf-8", errors="replace").strip()


# Extension → extractor mapping.
_EXTRACTORS: dict[str, Callable[[bytes], str]] = {
    "pdf": extract_text_from_pdf,
    "docx": extract_text_from_docx,
    "doc": extract_text_from_docx,
    "xlsx": extract_text_from_xlsx,
    "xls": extract_text_from_xlsx,
    "pptx": extract_text_from_pptx,
    "ppt": extract_text_from_pptx,
    "hwp": extract_text_from_hwp,
    "md": extract_text_from_md,
    "markdown": extract_text_from_md,
    "txt": extract_text_from_txt,
    "text": extract_text_from_txt,
    "csv": extract_text_from_txt,
}


def extract_text(data: bytes, ext: str) -> str:
    """Route to the appropriate extractor based on file extension.

    Args:
        data: Raw file bytes.
        ext: File extension without dot (e.g. "pdf", "docx").

    Returns:
        Extracted text or an error message.
    """
    extractor = _EXTRACTORS.get(ext.lower())
    if extractor is None:
        return f"(지원하지 않는 파일 형식이에요: .{ext})"
    return extractor(data)
